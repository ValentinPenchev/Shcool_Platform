import io
from pptx import Presentation


def _slide_xml(slide) -> str:
    """Суровият XML на слайда - най-надеждният начин да се засекат елементи
    (медия, SmartArt, преходи, анимации), за които python-pptx няма кратко API."""
    return slide.element.xml


def _has_media(prs) -> bool:
    return any("videoFile" in _slide_xml(s) or "audioFile" in _slide_xml(s) for s in prs.slides)


def _has_smartart(prs) -> bool:
    return any("drawingml/2006/diagram" in _slide_xml(s) for s in prs.slides)


def _has_transitions(prs) -> bool:
    return any("<p:transition" in _slide_xml(s) for s in prs.slides)


def _has_animations(prs) -> bool:
    # "<p:timing" (без затваряща ъглова скоба) улавя и самозатварящия се вариант "<p:timing/>"
    return any("<p:timing" in _slide_xml(s) for s in prs.slides)


def evaluate_ppt(file_bytes: bytes, criteria: dict) -> dict:
    """Оценява .pptx файл спрямо предоставените критерии."""
    prs = Presentation(io.BytesIO(file_bytes))

    score = 0
    max_score = 0
    details = []

    def add_check(key, label, passed_note, failed_note, passed):
        nonlocal score, max_score
        rule = criteria[key]
        rule_points = rule.get("points", 1) if isinstance(rule, dict) else 1
        max_score += rule_points
        if passed:
            score += rule_points
        details.append({
            "criterion": label,
            "passed": passed,
            "score": rule_points if passed else 0,
            "max": rule_points,
            "note": passed_note if passed else failed_note
        })

    # 1. Брой слайдове (min_slides)
    if "min_slides" in criteria:
        rule = criteria["min_slides"]
        required_slides = rule.get("count", 3) if isinstance(rule, dict) else 3
        slides_count = len(prs.slides)
        add_check(
            "min_slides", "Брой слайдове",
            f"Презентацията съдържа {slides_count} слайда (минимално изисквани: {required_slides}).",
            f"Презентацията съдържа само {slides_count} слайда (минимално изисквани: {required_slides}).",
            slides_count >= required_slides
        )

    # 2. Вмъкнато аудио или видео
    if "has_media" in criteria:
        add_check(
            "has_media", "Вмъкнато аудио/видео",
            "Намерено е вмъкнато аудио или видео в презентацията.",
            "Липсва вмъкнато аудио или видео.",
            _has_media(prs)
        )

    # 3. SmartArt графика
    if "has_smartart" in criteria:
        add_check(
            "has_smartart", "SmartArt графика",
            "Намерена е SmartArt графика в презентацията.",
            "Липсва SmartArt графика.",
            _has_smartart(prs)
        )

    # 4. Преходи между слайдовете
    if "has_transitions" in criteria:
        add_check(
            "has_transitions", "Преходи между слайдовете",
            "Намерен е зададен преход поне на един слайд.",
            "Не са зададени преходи между слайдовете.",
            _has_transitions(prs)
        )

    # 5. Анимации на обекти
    if "has_animations" in criteria:
        add_check(
            "has_animations", "Анимация на обекти",
            "Намерена е зададена анимация поне на един обект.",
            "Не са зададени анимации на обекти в слайдовете.",
            _has_animations(prs)
        )

    percentage = round((score / max_score) * 100) if max_score > 0 else 0
    return {
        "score": score,
        "max_score": max_score,
        "percentage": percentage,
        "details": details
    }
